#!/usr/bin/env python3
"""Generate the AWS Student Builder Group workshop deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree


# --- Brand ---------------------------------------------------------------
NAVY = RGBColor(0x12, 0x16, 0x20)
NAVY_RAISED = RGBColor(0x18, 0x1E, 0x2B)
NAVY_CARD = RGBColor(0x1B, 0x22, 0x32)
NAVY_LINE = RGBColor(0x2A, 0x33, 0x48)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
PURPLE_DEEP = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_SOFT = RGBColor(0xC4, 0xB5, 0xFD)
PURPLE_DIM = RGBColor(0x6D, 0x28, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MIST = RGBColor(0xE5, 0xE7, 0xEB)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)
SLATE = RGBColor(0x6B, 0x72, 0x80)

W = Inches(13.333)
H = Inches(7.5)
SLIDE_W = 13.333
SLIDE_H = 7.5

FONT_UI = "Inter"
FONT_MONO = "JetBrains Mono"


def _set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _disable_shadow(shape):
    sp_pr = shape._element.spPr
    effect = sp_pr.find(qn("a:effectLst"))
    if effect is not None:
        sp_pr.remove(effect)


def _set_run(paragraph, text, *, size, color, font=FONT_UI, bold=False):
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    rpr = run._r.get_or_add_rPr()
    latin = rpr.find(qn("a:latin"))
    if latin is None:
        latin = etree.SubElement(rpr, qn("a:latin"))
    latin.set("typeface", font)
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", font)
    return run


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=16,
    color=WHITE,
    font=FONT_UI,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    _set_run(p, text, size=size, color=color, font=font, bold=bold)
    return box


def add_rect(slide, left, top, width, height, color, *, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    _set_shape_fill(shape, color)
    _disable_shadow(shape)
    if radius is not None:
        # python-pptx adj is 0-1 relative to min(width,height)/2
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    return shape


def paint_background(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Subtle construction grid
    step = 0.42
    x = 0.0
    while x <= SLIDE_W + 0.01:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(0), Inches(0.007), Inches(SLIDE_H)
        )
        _set_shape_fill(line, NAVY_LINE)
        _disable_shadow(line)
        x += step
    y = 0.0
    while y <= SLIDE_H + 0.01:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(y), Inches(SLIDE_W), Inches(0.007)
        )
        _set_shape_fill(line, NAVY_LINE)
        _disable_shadow(line)
        y += step

    # Left brand rail
    add_rect(slide, 0, 0, 0.085, SLIDE_H, PURPLE)

    # Pixel-block accents (flyer language)
    blocks = [
        (12.55, 0.00, 0.78, 0.28, PURPLE),
        (12.95, 0.28, 0.38, 0.38, PURPLE_DEEP),
        (12.18, 0.00, 0.37, 0.18, PURPLE_DIM),
        (0.00, 7.12, 0.55, 0.38, PURPLE),
        (0.55, 7.28, 0.32, 0.22, PURPLE_DEEP),
        (12.78, 7.18, 0.55, 0.32, PURPLE),
        (12.40, 7.38, 0.28, 0.12, PURPLE_DIM),
    ]
    for left, top, width, height, color in blocks:
        add_rect(slide, left, top, width, height, color)


def add_footer(slide, index, total):
    add_rect(slide, 0.42, 7.12, 12.5, 0.012, NAVY_LINE)
    add_textbox(
        slide,
        0.42,
        7.16,
        8.6,
        0.28,
        "AWS STUDENT BUILDER GROUP  ·  JKUAT  ·  19 AUGUST 2026",
        size=9,
        color=SLATE,
        font=FONT_MONO,
        bold=False,
    )
    add_textbox(
        slide,
        10.55,
        7.16,
        2.3,
        0.28,
        f"{index:02d}  /  {total:02d}",
        size=9,
        color=PURPLE_SOFT,
        font=FONT_MONO,
        align=PP_ALIGN.RIGHT,
    )


def add_eyebrow(slide, label):
    add_rect(slide, 0.48, 0.32, 0.22, 0.08, PURPLE, radius=0.5)
    add_textbox(
        slide,
        0.80,
        0.24,
        8.5,
        0.28,
        label.upper(),
        size=11,
        color=PURPLE_SOFT,
        font=FONT_MONO,
        bold=True,
    )


def add_title(slide, title, subtitle=None):
    add_textbox(
        slide,
        0.48,
        0.52,
        12.2,
        0.58,
        title,
        size=30,
        color=WHITE,
        font=FONT_UI,
        bold=True,
    )
    if subtitle:
        add_textbox(
            slide,
            0.48,
            1.08,
            12.2,
            0.36,
            subtitle,
            size=14,
            color=MUTED,
            font=FONT_UI,
        )


def add_sentence_cards(slide, sentences, *, top=1.62, card_h=2.20, row_gap=0.22):
    """Four rounded cards in a 2x2 grid — one sentence each."""
    positions = [
        (0.48, top),
        (6.92, top),
        (0.48, top + card_h + row_gap),
        (6.92, top + card_h + row_gap),
    ]
    card_w = 5.92
    for i, ((left, card_top), sentence) in enumerate(zip(positions, sentences), start=1):
        add_rect(slide, left, card_top, card_w, card_h, NAVY_CARD, radius=0.08)
        add_rect(slide, left, card_top, 0.09, card_h, PURPLE)
        add_textbox(
            slide,
            left + 0.28,
            card_top + 0.18,
            1.4,
            0.26,
            f"0{i}",
            size=12,
            color=PURPLE,
            font=FONT_MONO,
            bold=True,
        )
        add_textbox(
            slide,
            left + 0.28,
            card_top + 0.50,
            card_w - 0.56,
            card_h - 0.68,
            sentence,
            size=15,
            color=MIST,
            font=FONT_UI,
        )


def add_demo_steps(slide, sentences):
    """Four wide numbered steps for the live demo."""
    top = 1.62
    for i, sentence in enumerate(sentences, start=1):
        y = top + (i - 1) * 1.28
        add_rect(slide, 0.48, y, 12.36, 1.16, NAVY_CARD, radius=0.08)
        add_rect(slide, 0.48, y, 1.16, 1.16, PURPLE_DEEP, radius=0.08)
        add_rect(slide, 0.48, y, 0.10, 1.16, PURPLE)
        add_textbox(
            slide,
            0.48,
            y + 0.30,
            1.16,
            0.56,
            f"0{i}",
            size=22,
            color=WHITE,
            font=FONT_MONO,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            1.90,
            y + 0.28,
            10.60,
            0.64,
            sentence,
            size=16,
            color=MIST,
            font=FONT_UI,
        )


def add_command_steps(slide, steps):
    """Four rows: one spoken sentence plus the command or console click."""
    top = 1.58
    for i, step in enumerate(steps, start=1):
        y = top + (i - 1) * 1.30
        add_rect(slide, 0.48, y, 12.36, 1.18, NAVY_CARD, radius=0.08)
        add_rect(slide, 0.48, y, 1.16, 1.18, PURPLE_DEEP, radius=0.08)
        add_rect(slide, 0.48, y, 0.10, 1.18, PURPLE)
        add_textbox(
            slide,
            0.48,
            y + 0.32,
            1.16,
            0.56,
            f"0{i}",
            size=20,
            color=WHITE,
            font=FONT_MONO,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            1.86,
            y + 0.12,
            10.70,
            0.36,
            step["label"],
            size=13,
            color=PURPLE_SOFT,
            font=FONT_UI,
        )
        add_textbox(
            slide,
            1.86,
            y + 0.48,
            10.70,
            0.58,
            step["command"],
            size=13,
            color=WHITE,
            font=FONT_MONO,
        )


def add_social_cards(slide, items):
    """Four contact cards — handle + one sentence each."""
    positions = [
        (0.48, 1.62),
        (6.92, 1.62),
        (0.48, 4.04),
        (6.92, 4.04),
    ]
    card_w, card_h = 5.92, 2.20
    for (left, top), item in zip(positions, items):
        add_rect(slide, left, top, card_w, card_h, NAVY_CARD, radius=0.08)
        add_rect(slide, left, top, card_w, 0.08, PURPLE)
        add_textbox(
            slide,
            left + 0.32,
            top + 0.28,
            card_w - 0.64,
            0.28,
            item["label"],
            size=11,
            color=PURPLE_SOFT,
            font=FONT_MONO,
            bold=True,
        )
        handle_box = add_textbox(
            slide,
            left + 0.32,
            top + 0.56,
            card_w - 0.64,
            0.40,
            item["handle"],
            size=18,
            color=WHITE,
            font=FONT_UI,
            bold=True,
        )
        if item.get("url"):
            run = handle_box.text_frame.paragraphs[0].runs[0]
            run.hyperlink.address = item["url"]
        add_textbox(
            slide,
            left + 0.32,
            top + 1.08,
            card_w - 0.64,
            0.88,
            item["sentence"],
            size=14,
            color=MIST,
            font=FONT_UI,
        )


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    paint_background(slide)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides_spec = []

    # 1 — Introduction
    slides_spec.append(
        {
            "kind": "intro",
            "eyebrow": "Introduction",
            "title": "From Backend Code to the Cloud",
            "subtitle": "Building and Deploying Production-Ready Apps on AWS",
            "sentences": [
                "Welcome to this AWS Student Builder Group session, where we turn local backend code into a service that real users can reach.",
                "I am Gilbert Chris, Founder at Remoteflow, and I will walk you from a working laptop API to a production-minded AWS deploy.",
                "Tonight is for student builders who already write routes and queries, and now want those apps to survive traffic, secrets, and failure.",
                "Stay with the flow, follow the demo, and leave knowing the exact next ship you can make after this workshop.",
            ],
        }
    )

    # 2
    slides_spec.append(
        {
            "kind": "cards",
            "eyebrow": "The local starting point",
            "title": "What a modern backend really is",
            "sentences": [
                "A modern backend is more than a framework folder that answers on localhost:8000.",
                "Production code must keep secrets out of git, recover from crashes, and tell you when something breaks.",
                "Many student projects stop at it works on my machine, and that gap is exactly what we close tonight.",
                "Treat your API as a product with identity, storage, logs, and a repeatable path from commit to cloud.",
            ],
        }
    )

    # 3
    slides_spec.append(
        {
            "kind": "cards",
            "eyebrow": "The platform choice",
            "title": "Why we deploy this on AWS",
            "sentences": [
                "AWS gives student builders the same compute, data, and security primitives that real companies compose in production.",
                "You stop babysitting one fragile VPS and start assembling services that scale, fail independently, and can be replaced.",
                "A generous free tier plus global regions means you can learn production thinking without needing a data center.",
                "Choose AWS when you want elasticity, managed databases, and a career-shaped skill that still fits a student budget.",
            ],
        }
    )

    # 4
    slides_spec.append(
        {
            "kind": "cards",
            "eyebrow": "The map",
            "title": "Architecture we will ship tonight",
            "sentences": [
                "Clients reach a public HTTPS load balancer, which forwards traffic to our container.",
                "The application runs on Amazon ECS Express Mode on Fargate, so we never SSH into a snowflake server.",
                "State lives in Amazon RDS, and every secret comes from Secrets Manager instead of a committed .env file.",
                "CloudWatch collects logs and metrics so we can see errors before a classmate in the chat reports them.",
            ],
        }
    )

    # 5 — Demo map
    slides_spec.append(
        {
            "kind": "demo",
            "eyebrow": "Live demo",
            "title": "Console first, laptop only for code",
            "sentences": [
                "Split the screen: PowerShell on the left for your API, and the AWS Console on the right for every cloud service.",
                "Stay in N. Virginia, zoom the browser, and use curl.exe on Windows because curl is not real curl.",
                "We pack the same FastAPI app in Docker, store it in ECR, then click Express Mode, Secrets Manager, RDS, and CloudWatch.",
                "When /health says cloud and /db returns ping 1, the identical API is live in AWS.",
            ],
        }
    )

    # 6 — Local + Docker commands
    slides_spec.append(
        {
            "kind": "commands",
            "eyebrow": "Beats 1 and 2",
            "title": "Laptop and Docker — PowerShell only",
            "steps": [
                {
                    "label": "Start the API from talks/2026-08-19-backend-to-cloud/demo.",
                    "command": "python -m pip install -r requirements.txt   then   python -m uvicorn app:app --reload --host 0.0.0.0 --port 8000",
                },
                {
                    "label": "Prove localhost works; /health must show stage local.",
                    "command": "curl.exe http://localhost:8000/health    curl.exe http://localhost:8000/    curl.exe -X POST http://localhost:8000/echo -H \"Content-Type: application/json\" -d \"{\\\"message\\\":\\\"JKUAT to the cloud\\\"}\"",
                },
                {
                    "label": "Stop uvicorn, then wrap the same app in a container.",
                    "command": "docker build -t workshop-api .    &&    docker run --rm -p 8000:8000 workshop-api",
                },
                {
                    "label": "Repeat the three curl.exe calls against localhost; the JSON must match.",
                    "command": "If Docker _ping 500: start Docker Desktop, wait for Engine running, open a new PowerShell.",
                },
            ],
        }
    )

    # 7 — ECR
    slides_spec.append(
        {
            "kind": "commands",
            "eyebrow": "Beat 3 · AWS Console",
            "title": "ECR — store the Docker box",
            "steps": [
                {
                    "label": "Search ECR, then create a private repository named workshop-api.",
                    "command": "Console: Elastic Container Registry  →  Repositories  →  Create repository  →  Private  →  workshop-api",
                },
                {
                    "label": "Open View push commands and choose the Windows tab.",
                    "command": "aws ecr get-login-password --region us-east-1 | docker login --username AWS --password-stdin ACCOUNT.dkr.ecr.us-east-1.amazonaws.com",
                },
                {
                    "label": "Tag and push the image you already built.",
                    "command": "docker tag workshop-api:latest ACCOUNT.dkr.ecr.us-east-1.amazonaws.com/workshop-api:latest    &&    docker push .../workshop-api:latest",
                },
                {
                    "label": "Refresh the repository until tag latest appears.",
                    "command": "ECR is the warehouse; ECS can only run a box that is stored here.",
                },
            ],
        }
    )

    # 8 — Secrets + Express Mode
    slides_spec.append(
        {
            "kind": "commands",
            "eyebrow": "Beats 4 and 5 · AWS Console",
            "title": "Secrets Manager and ECS Express Mode",
            "steps": [
                {
                    "label": "Store a secret named workshop/api; do not retrieve the value on the projector.",
                    "command": "Console: Secrets Manager  →  Store a new secret  →  Other type  →  name workshop/api",
                },
                {
                    "label": "Create the service from the ECR image; let the console create both IAM roles.",
                    "command": "Console: ECS  →  Express mode  →  Create  →  Browse ECR  →  workshop-api  →  latest  →  Create new role × 2",
                },
                {
                    "label": "Open Additional configurations and set the port, health check, and stage.",
                    "command": "Name workshop-api   ·   Container port 8000   ·   Health check /health   ·   APP_STAGE=cloud   ·   PORT=8000",
                },
                {
                    "label": "Wait until Active, then copy the Application URL.",
                    "command": "https://workshop-api.ecs.us-east-1.on.aws    —    narrate ALB, Fargate, security groups, CloudWatch on the Timeline",
                },
            ],
        }
    )

    # 9 — Live test
    slides_spec.append(
        {
            "kind": "commands",
            "eyebrow": "Beat 6",
            "title": "Same curls, now on the cloud URL",
            "steps": [
                {
                    "label": "Paste your Application URL into PowerShell.",
                    "command": "$App = \"https://workshop-api.ecs.us-east-1.on.aws\"",
                },
                {
                    "label": "Replay the three requests; /health must show stage cloud.",
                    "command": "curl.exe $App/health    curl.exe $App/    curl.exe -X POST $App/echo -H \"Content-Type: application/json\" -d \"{\\\"message\\\":\\\"JKUAT to the cloud\\\"}\"",
                },
                {
                    "label": "Open CloudWatch Logs and refresh after another /health.",
                    "command": "Console: CloudWatch  →  Logs  →  Log groups  →  newest stream  →  look for GET /health stage=cloud",
                },
                {
                    "label": "That is the punchline: identical routes, different place.",
                    "command": "If health is unhealthy, the container port is still 80 — change it to 8000 and redeploy.",
                },
            ],
        }
    )

    # 10 — RDS console
    slides_spec.append(
        {
            "kind": "commands",
            "eyebrow": "Beat 7 · AWS Console",
            "title": "RDS PostgreSQL — the filing cabinet",
            "steps": [
                {
                    "label": "Create Postgres in the default VPC; wait until Available.",
                    "command": "RDS  →  Create database  →  PostgreSQL  →  workshop-db  →  user workshop  →  Public access No  →  SG workshop-db-sg",
                },
                {
                    "label": "Allow only the ECS task security group into port 5432.",
                    "command": "EC2  →  Security Groups  →  workshop-db-sg  →  Inbound  →  PostgreSQL 5432  →  source = task sg-...   never 0.0.0.0/0",
                },
                {
                    "label": "Put the connection string in Secrets Manager, not in Git.",
                    "command": "Secret workshop/database-url  =  postgresql://workshop:PASSWORD@ENDPOINT:5432/workshop",
                },
                {
                    "label": "Update Express Mode and inject the secret as DATABASE_URL.",
                    "command": "Express mode  →  Update  →  Environment variables  →  DATABASE_URL  →  Value type Secret  →  workshop/database-url",
                },
            ],
        }
    )

    # 11 — Test RDS
    slides_spec.append(
        {
            "kind": "commands",
            "eyebrow": "Beat 7 · prove it",
            "title": "How we test that RDS is working",
            "steps": [
                {
                    "label": "RDS console must show Available and an endpoint before any curl.",
                    "command": "RDS  →  workshop-db  →  Connectivity  →  copy Endpoint   ·   status badge Available",
                },
                {
                    "label": "/health reports whether the task has a working connection.",
                    "command": "curl.exe $App/health      expect   \"database\":\"connected\"",
                },
                {
                    "label": "GET /db is the real proof: the API runs SELECT 1 inside Postgres.",
                    "command": "curl.exe $App/db      expect   {\"status\":\"connected\",\"ping\":1,\"database_name\":\"workshop\",\"engine\":\"PostgreSQL\"}",
                },
                {
                    "label": "Confirm in logs and on the RDS graph so beginners see it two more ways.",
                    "command": "CloudWatch: GET /db status=connected     ·     RDS → Monitoring → Database connections leaves zero",
                },
            ],
        }
    )

    # 12
    slides_spec.append(
        {
            "kind": "cards",
            "eyebrow": "After the first deploy",
            "title": "Habits that make it production-ready",
            "sentences": [
                "Never commit access keys, and rotate anything that has already appeared in a screenshot or a chat paste.",
                "Add health checks, structured logs, and one real alarm before you invite classmates to traffic-test the URL.",
                "Ship through CI/CD so every change is tested, reviewed, and easy to roll back when a release misbehaves.",
                "Least-privilege IAM, HTTPS everywhere, and automated backups are the price of calling an app production-ready.",
            ],
        }
    )

    # 7
    slides_spec.append(
        {
            "kind": "cards",
            "eyebrow": "Before we close",
            "title": "What you should build this week",
            "sentences": [
                "Redeploy tonight's service under your own AWS account before the energy of the session fades.",
                "Add one production habit immediately: move secrets, add logging, or wire a tiny CI pipeline.",
                "Share the public URL in the AWS Student Builder Group so another builder can learn from your path.",
                "If a step blocks you, ask in the community; shipped questions teach more than perfect private repos.",
            ],
        }
    )

    # 8 — Socials
    slides_spec.append(
        {
            "kind": "socials",
            "eyebrow": "Socials",
            "title": "Stay in touch and keep shipping",
            "items": [
                {
                    "label": "LINKEDIN",
                    "handle": "Gilbert Chris",
                    "url": "https://www.linkedin.com/in/gilbert-chris-a696151a9",
                    "sentence": "Connect with me on LinkedIn as Gilbert Chris for recaps, AWS notes, and what we build next at Remoteflow.",
                },
                {
                    "label": "GITHUB",
                    "handle": "github.com/GilbertKamau",
                    "url": "https://github.com/GilbertKamau",
                    "sentence": "Follow the public repos on GitHub at GilbertKamau and fork the demo path we used from laptop to cloud.",
                },
                {
                    "label": "WEB",
                    "handle": "remoteflow.cc",
                    "url": "https://www.remoteflow.cc/",
                    "sentence": "Visit remoteflow.cc to see a production product born from the same backend-to-cloud discipline we practiced tonight.",
                },
                {
                    "label": "X  ·  EMAIL",
                    "handle": "@El_berto31",
                    "url": "https://x.com/El_berto31",
                    "sentence": "Find me on X as @El_berto31, or write gilbertchris062@gmail.com if you want feedback on the service you ship this week.",
                },
            ],
        }
    )

    total = len(slides_spec)

    for index, spec in enumerate(slides_spec, start=1):
        slide = new_slide(prs)
        kind = spec["kind"]

        if kind == "intro":
            add_eyebrow(slide, spec["eyebrow"])
            add_textbox(
                slide,
                0.48,
                0.50,
                12.2,
                0.70,
                spec["title"],
                size=32,
                color=WHITE,
                font=FONT_UI,
                bold=True,
            )
            add_textbox(
                slide,
                0.48,
                1.18,
                12.2,
                0.34,
                spec["subtitle"],
                size=16,
                color=PURPLE_SOFT,
                font=FONT_UI,
            )

            # Event chips
            chips = [
                (0.48, "WED 19 AUG 2026"),
                (3.05, "07:30 – 09:30 PM"),
                (5.62, "ONLINE"),
                (7.42, "GILBERT CHRIS  ·  REMOTEFLOW"),
            ]
            widths = [2.42, 2.42, 1.65, 5.40]
            for (left, label), width in zip(chips, widths):
                add_rect(slide, left, 1.64, width, 0.38, NAVY_CARD, radius=0.4)
                add_textbox(
                    slide,
                    left,
                    1.70,
                    width,
                    0.28,
                    label,
                    size=10,
                    color=MIST,
                    font=FONT_MONO,
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )

            add_sentence_cards(slide, spec["sentences"], top=2.18, card_h=2.14, row_gap=0.18)

        elif kind == "cards":
            add_eyebrow(slide, spec["eyebrow"])
            add_title(slide, spec["title"])
            add_sentence_cards(slide, spec["sentences"])

        elif kind == "demo":
            add_eyebrow(slide, spec["eyebrow"])
            add_title(slide, spec["title"])
            add_demo_steps(slide, spec["sentences"])
            notes = spec.get("notes")
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        elif kind == "commands":
            add_eyebrow(slide, spec["eyebrow"])
            add_title(slide, spec["title"])
            add_command_steps(slide, spec["steps"])

        elif kind == "socials":
            add_eyebrow(slide, spec["eyebrow"])
            add_title(slide, spec["title"])
            add_social_cards(slide, spec["items"])

        add_footer(slide, index, total)

    out_dir = Path(__file__).resolve().parent
    out_path = out_dir / "From-Backend-Code-to-the-Cloud.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")
    return out_path


if __name__ == "__main__":
    build()
